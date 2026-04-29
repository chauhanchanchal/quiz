# AI-Powered Quiz & Exam Platform
from flask import Flask, render_template, request, redirect, session, jsonify, Response, send_file, url_for


import json, os, io, time, requests as http_requests, datetime, csv, random, copy, uuid, time
import PyPDF2
from pptx import Presentation
import ppt2txt, tempfile, zipfile
try:
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False


app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "secret123")

DATA_FILE = "data.json"
GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "")

# ---------- DATA ----------
def load_data():
    if not os.path.exists(DATA_FILE):
        return {"users": [], "quizzes": [], "notes": [], "results": []}
    with open(DATA_FILE, "r") as f:
        data = json.load(f)
    for key in ["quizzes", "notes", "results", "answer_reports"]:
        if key not in data:
            data[key] = []
    return data

def save_data(data):
    with open(DATA_FILE, "w") as f:
        json.dump(data, f, indent=4)

# ---------- BADGES ----------
def compute_badges(student_id, results):
    student_results = [r for r in results if r["student_id"] == student_id]
    badges = []
    if not student_results:
        return badges
    scores = [(r["score"] / r["total"] * 100) if r["total"] > 0 else 0 for r in student_results]
    if any(s == 100 for s in scores):
        badges.append({"icon": "🏆", "label": "Perfect Score", "color": "#f59e0b"})
    if len(student_results) >= 1:
        badges.append({"icon": "🎯", "label": "First Attempt", "color": "#3b82f6"})
    if len(student_results) >= 5:
        badges.append({"icon": "🔥", "label": "Quiz Streak", "color": "#ef4444"})
    if len(student_results) >= 10:
        badges.append({"icon": "⭐", "label": "Quiz Master", "color": "#6366f1"})
    if sum(1 for s in scores if s >= 70) >= 3:
        badges.append({"icon": "📈", "label": "High Achiever", "color": "#4ade80"})
    if sum(1 for s in scores if s == 100) >= 3:
        badges.append({"icon": "💎", "label": "Perfectionist", "color": "#a78bfa"})
    return badges

# ---------- LOGIN ----------
@app.route("/", methods=["GET", "POST"])
def login():
    if request.method == "POST":
        ku_id = request.form["ku_id"]
        password = request.form["password"]
        data = load_data()
        for user in data["users"]:
            if user["ku_id"] == ku_id and user["password"] == password:
                session["user"] = user
                return redirect("/mentor" if user["role"] == "mentor" else "/student")
        return render_template("login.html", error="Invalid Credentials")
    return render_template("login.html")

# ---------- REGISTER ----------
@app.route("/register", methods=["GET", "POST"])
def register():
    if request.method == "POST":
        data = load_data()
        new_user = {
            "name": request.form["name"], "ku_id": request.form["ku_id"],
            "password": request.form["password"], "class": request.form["class"], "role": "student"
        }
        for u in data["users"]:
            if u["ku_id"] == new_user["ku_id"]:
                return render_template("register.html", error="User already exists")
        data["users"].append(new_user)
        save_data(data)
        session["user"] = new_user
        return redirect("/student")
    return render_template("register.html")

# ---------- LOGOUT ----------
@app.route("/logout")
def logout():
    session.pop("user", None)
    return redirect("/")

# ---------- STUDENT DASHBOARD ----------
@app.route("/student")
def student():
    if "user" not in session:
        return redirect("/")
    user = session["user"]
    data = load_data()
    now = datetime.datetime.now()

    # Count attempts per quiz for this student
    attempt_counts = {}
    best_scores = {}
    for r in data.get("results", []):
        if r["student_id"] == user["ku_id"]:
            qid = r["quiz_id"]
            attempt_counts[qid] = attempt_counts.get(qid, 0) + 1
            pct = round(r["score"] / r["total"] * 100, 1) if r["total"] > 0 else 0
            if qid not in best_scores or pct > best_scores[qid]:
                best_scores[qid] = pct

    student_quizzes = []
    for i, q in enumerate(data["quizzes"]):
        if q["class"] not in [user["class"], "All"]:
            continue
        open_at = q.get("open_at", "")
        close_at = q.get("close_at", "")
        if open_at:
            try:
                if datetime.datetime.strptime(open_at, "%Y-%m-%dT%H:%M") > now:
                    continue
            except:
                pass
        if close_at:
            try:
                if datetime.datetime.strptime(close_at, "%Y-%m-%dT%H:%M") < now:
                    continue
            except:
                pass
        quiz_data = dict(q, id=i)
        quiz_data["attempts_used"] = attempt_counts.get(i, 0)
        quiz_data["max_attempts"] = q.get("max_attempts", 0)
        quiz_data["best_score"] = best_scores.get(i, None)
        quiz_data["is_locked"] = (quiz_data["max_attempts"] > 0 and quiz_data["attempts_used"] >= quiz_data["max_attempts"])
        student_quizzes.append(quiz_data)

    results = [r for r in data.get("results", []) if r["student_id"] == user["ku_id"]]
    results = sorted(results, key=lambda x: x.get("timestamp", ""), reverse=True)
    badges = compute_badges(user["ku_id"], data.get("results", []))

    # --- Analytics computation ---
    all_results = data.get("results", [])
    my_results = [r for r in all_results if r["student_id"] == user["ku_id"]]

    # Score history for line chart
    score_history = []
    for r in sorted(my_results, key=lambda x: x.get("timestamp", "")):
        pct = round(r["score"] / r["total"] * 100, 1) if r["total"] > 0 else 0
        score_history.append({"quiz": r["quiz_title"], "pct": pct, "ts": r.get("timestamp", "")})

    # Subject-wise stats for radar chart
    subject_stats = {}
    for r in my_results:
        subject = "General"
        if 0 <= r["quiz_id"] < len(data["quizzes"]):
            subject = data["quizzes"][r["quiz_id"]].get("subject", "General")
        if subject not in subject_stats:
            subject_stats[subject] = {"scores": [], "total_attempts": 0}
        pct = round(r["score"] / r["total"] * 100, 1) if r["total"] > 0 else 0
        subject_stats[subject]["scores"].append(pct)
        subject_stats[subject]["total_attempts"] += 1
    for sub in subject_stats:
        scores = subject_stats[sub]["scores"]
        subject_stats[sub]["avg"] = round(sum(scores) / len(scores), 1) if scores else 0

    # Strengths and weaknesses
    sorted_subjects = sorted(subject_stats.items(), key=lambda x: x[1]["avg"], reverse=True)
    strengths = [s[0] for s in sorted_subjects if s[1]["avg"] >= 60][:3]
    weaknesses = [s[0] for s in sorted_subjects if s[1]["avg"] < 60][:3]

    # Percentile calculation
    student_averages = {}
    for r in all_results:
        sid = r["student_id"]
        if sid not in student_averages:
            student_averages[sid] = []
        pct = round(r["score"] / r["total"] * 100, 1) if r["total"] > 0 else 0
        student_averages[sid].append(pct)
    my_avg = 0
    percentile = 0
    if user["ku_id"] in student_averages:
        my_scores = student_averages[user["ku_id"]]
        my_avg = round(sum(my_scores) / len(my_scores), 1)
        all_avgs = [round(sum(v) / len(v), 1) for v in student_averages.values()]
        below_count = sum(1 for a in all_avgs if a < my_avg)
        percentile = round(below_count / len(all_avgs) * 100) if all_avgs else 0

    analytics = {
        "score_history": score_history,
        "subject_stats": subject_stats,
        "strengths": strengths,
        "weaknesses": weaknesses,
        "percentile": percentile,
        "my_avg": my_avg,
        "total_quizzes": len(my_results)
    }

    return render_template("student_dashboard.html", user=user, quizzes=student_quizzes,
                         results=results, badges=badges, analytics=analytics)

# ---------- MENTOR DASHBOARD ----------
@app.route("/mentor")
def mentor():
    if "user" not in session or session["user"].get("role") != "mentor":
        return redirect("/")
    data = load_data()
    
    mentor_quizzes = [dict(q, id=i) for i, q in enumerate(data["quizzes"])]
    results = sorted(data.get("results", []), key=lambda x: x.get("timestamp", ""), reverse=True)
    
    # 1. Class Analytics per quiz
    for q in mentor_quizzes:
        q_results = [r for r in results if r["quiz_id"] == q["id"]]
        if q_results:
            q["avg_score"] = round(sum(r["score"]/r["total"]*100 for r in q_results if r["total"]>0) / len(q_results), 1)
            top = max(q_results, key=lambda x: x["score"]/x["total"] if x["total"]>0 else 0)
            q["top_performer"] = top["student_name"]
            q["attempts_count"] = len(q_results)
        else:
            q["avg_score"] = 0
            q["top_performer"] = "N/A"
            q["attempts_count"] = 0

    # 3. Student Roster stats
    students = [u for u in data["users"] if u.get("role") == "student"]
    roster = []
    for s in students:
        s_results = [r for r in results if r["student_id"] == s["ku_id"]]
        if s_results:
            avg = round(sum(r["score"]/r["total"]*100 for r in s_results if r["total"]>0) / len(s_results), 1)
        else:
            avg = 0
            
        assigned = sum(1 for q in data["quizzes"] if q.get("class") in ["All", s.get("class")])
        taken_quizzes = len(set(r["quiz_id"] for r in s_results))
        missed = assigned - taken_quizzes
        
        roster.append({
            "name": s["name"],
            "ku_id": s["ku_id"],
            "class": s.get("class", "N/A"),
            "avg_score": avg,
            "total_taken": len(s_results),
            "missed": max(0, missed)
        })
        
    roster.sort(key=lambda x: x["avg_score"], reverse=True)

    answer_reports = sorted(data.get("answer_reports", []), key=lambda x: x.get("submitted_at",""), reverse=True)
    return render_template("mentor_dashboard.html", user=session["user"], quizzes=mentor_quizzes, results=results, roster=roster, answer_reports=answer_reports)

# ---------- EXPORT CSV ----------
@app.route("/export-csv/<int:quiz_id>")
def export_csv(quiz_id):
    if "user" not in session or session["user"].get("role") != "mentor":
        return redirect("/")
    data = load_data()
    if quiz_id < 0 or quiz_id >= len(data["quizzes"]):
        return redirect("/mentor")
        
    quiz_results = [r for r in data.get("results", []) if r["quiz_id"] == quiz_id]
    
    def generate():
        yield "Student ID,Student Name,Score,Total,Percentage,Timestamp\n"
        for r in quiz_results:
            pct = round(r["score"] / r["total"] * 100, 1) if r["total"] > 0 else 0
            ts = r.get("timestamp", "").replace(",", "")
            yield f"{r['student_id']},{r['student_name']},{r['score']},{r['total']},{pct}%,{ts}\n"
            
    return Response(generate(), mimetype="text/csv", headers={"Content-Disposition": f"attachment; filename=quiz_{quiz_id}_results.csv"})

# ---------- CREATE QUIZ ----------
@app.route("/create_quiz", methods=["POST"])
def create_quiz():
    if "user" not in session or session["user"].get("role") != "mentor":
        return redirect("/")
    data = load_data()
    if request.is_json:
        payload = request.get_json()
        new_quiz = {
            "title": payload.get("title", "Untitled Quiz"), "subject": payload.get("subject", "General"),
            "type": payload.get("type", "MCQ"), "class": payload.get("class", "All"),
            "questions": payload.get("questions", []), "created_by": session["user"]["ku_id"],
            "open_at": payload.get("open_at", ""), "close_at": payload.get("close_at", ""),
            "max_attempts": int(payload.get("max_attempts", 0)),
            "shuffle_questions": payload.get("shuffle_questions", False),
            "shuffle_options": payload.get("shuffle_options", False),
            "rapid_fire": payload.get("rapid_fire", False)
        }
        data["quizzes"].append(new_quiz)
        save_data(data)
        return jsonify({"success": True})

    try:
        questions = json.loads(request.form.get("questions_json", "[]"))
    except:
        questions = []

    new_quiz = {
        "title": request.form.get("title", "Untitled"), "subject": request.form.get("subject", "General"),
        "type": request.form.get("type", "MCQ"), "class": request.form.get("class", "All"),
        "questions": questions, "created_by": session["user"]["ku_id"],
        "open_at": request.form.get("open_at", ""), "close_at": request.form.get("close_at", ""),
        "max_attempts": int(request.form.get("max_attempts", 0)),
        "shuffle_questions": request.form.get("shuffle_questions") == "on",
        "shuffle_options": request.form.get("shuffle_options") == "on",
        "rapid_fire": request.form.get("rapid_fire") == "on"
    }
    data["quizzes"].append(new_quiz)
    save_data(data)
    return redirect("/mentor")

# ---------- EDIT QUIZ ----------
@app.route("/edit_quiz/<int:quiz_id>", methods=["GET", "POST"])
def edit_quiz(quiz_id):
    if "user" not in session or session["user"]["role"] != "mentor":
        return redirect("/")
    data = load_data()
    if quiz_id < 0 or quiz_id >= len(data["quizzes"]):
        return redirect("/mentor")

    if request.method == "POST":
        quiz = data["quizzes"][quiz_id]
        quiz["title"] = request.form.get("title", quiz["title"])
        quiz["subject"] = request.form.get("subject", quiz["subject"])
        quiz["type"] = request.form.get("type", quiz["type"])
        quiz["class"] = request.form.get("class", quiz["class"])
        quiz["open_at"] = request.form.get("open_at", "")
        quiz["close_at"] = request.form.get("close_at", "")
        quiz["max_attempts"] = int(request.form.get("max_attempts", 0))
        quiz["shuffle_questions"] = request.form.get("shuffle_questions") == "on"
        quiz["shuffle_options"] = request.form.get("shuffle_options") == "on"
        quiz["rapid_fire"] = request.form.get("rapid_fire") == "on"
        try:
            quiz["questions"] = json.loads(request.form.get("questions_json", "[]"))
        except:
            pass
        save_data(data)
        return redirect("/mentor")

    quiz = dict(data["quizzes"][quiz_id], id=quiz_id)
    return render_template("edit_quiz.html", quiz=quiz, user=session["user"])

# ---------- DELETE QUIZ ----------
@app.route("/delete_quiz/<int:quiz_id>")
def delete_quiz(quiz_id):
    if "user" not in session or session["user"]["role"] != "mentor":
        return redirect("/")
    data = load_data()
    if 0 <= quiz_id < len(data["quizzes"]):
        data["quizzes"].pop(quiz_id)
        new_results = []
        for r in data.get("results", []):
            if r["quiz_id"] == quiz_id:
                continue
            if r["quiz_id"] > quiz_id:
                r["quiz_id"] -= 1
            new_results.append(r)
        data["results"] = new_results
        save_data(data)
    return redirect("/mentor")

# ---------- EXPORT CSV ----------
@app.route("/export_results")
def export_results():
    if "user" not in session or session["user"]["role"] != "mentor":
        return redirect("/")
    data = load_data()
    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow(["Student Name", "KU ID", "Class", "Quiz Title", "Score", "Total", "Percentage", "Timestamp"])
    for r in data.get("results", []):
        pct = round(r["score"] / r["total"] * 100, 1) if r["total"] > 0 else 0
        writer.writerow([r["student_name"], r["student_id"], r.get("student_class", ""), r["quiz_title"], r["score"], r["total"], f"{pct}%", r["timestamp"]])
    return Response(output.getvalue(), mimetype="text/csv",
                    headers={"Content-Disposition": "attachment; filename=quiz_results.csv"})

# ---------- TAKE QUIZ ----------
@app.route("/quiz/<int:quiz_id>")
def take_quiz(quiz_id):
    if "user" not in session:
        return redirect("/")
    data = load_data()
    if quiz_id < 0 or quiz_id >= len(data["quizzes"]):
        return redirect("/student")
    quiz = dict(data["quizzes"][quiz_id], id=quiz_id)
    user = session["user"]

    # Check attempt limits
    max_attempts = quiz.get("max_attempts", 0)
    if max_attempts > 0:
        attempts_used = sum(1 for r in data.get("results", [])
                          if r["student_id"] == user["ku_id"] and r["quiz_id"] == quiz_id)
        if attempts_used >= max_attempts:
            return redirect("/student")

    # Shuffle questions if enabled
    questions = copy.deepcopy(quiz["questions"])
    if quiz.get("shuffle_questions", False):
        random.shuffle(questions)

    # Shuffle MCQ options if enabled
    if quiz.get("shuffle_options", False):
        for q in questions:
            if q.get("options"):
                random.shuffle(q["options"])

    quiz["questions"] = questions
    return render_template("quiz.html", quiz=quiz)

# ---------- SUBMIT QUIZ ----------
@app.route("/submit/<int:quiz_id>", methods=["POST"])
def submit_quiz(quiz_id):
    if "user" not in session:
        return redirect("/")
    data = load_data()
    if quiz_id < 0 or quiz_id >= len(data["quizzes"]):
        return redirect("/student")
    quiz = data["quizzes"][quiz_id]
    questions = quiz["questions"]
    score = 0
    total = len(questions)
    feedback = []

    for i, q in enumerate(questions):
        user_answer = request.form.get(str(i), "").strip()
        correct = q.get("answer", "").strip()
        options = q.get("options", [])
        correct_text = correct
        if len(correct) == 1 and correct.upper() in ["A", "B", "C", "D"] and options:
            idx = ord(correct.upper()) - ord("A")
            if idx < len(options):
                correct_text = options[idx].strip()
        is_correct = user_answer.lower() == correct_text.lower()
        if not is_correct:
            cu = "".join(c for c in user_answer.lower() if c.isalnum())
            cc = "".join(c for c in correct_text.lower() if c.isalnum())
            if cu == cc and cc:
                is_correct = True
        if quiz.get("type") == "Coding":
            is_correct = False
            explanation = f"Reference Solution: {correct_text}"
        else:
            explanation = "Correct!" if is_correct else f"The correct answer is: {correct_text}"
        if is_correct and quiz.get("type") != "Coding":
            score += 1
        feedback.append({"question": q["question"], "your": user_answer or "(no answer)", "correct": correct_text, "is_correct": is_correct, "explanation": explanation})

    user = session["user"]
    if "results" not in data:
        data["results"] = []
    data["results"].append({
        "student_id": user["ku_id"], "student_name": user["name"],
        "student_class": user.get("class", "N/A"), "quiz_id": quiz_id,
        "quiz_title": quiz["title"], "score": score, "total": total,
        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    })
    save_data(data)
    badges = compute_badges(user["ku_id"], data["results"])
    percentage = round(score / total * 100, 1) if total > 0 else 0
    return render_template("result.html", score=score, total=total, feedback=feedback, badges=badges, percentage=percentage)

# ---------- NOTES ----------
@app.route("/notes", methods=["GET", "POST"])
def notes():
    if "user" not in session:
        return redirect("/")

    data = load_data()
    user = session["user"]
    user_id = user["ku_id"]

    if request.method == "POST":
        note_text = request.form.get("note", "").strip()
        visibility = request.form.get("visibility", "public")
        target_id = request.form.get("target_id")
        target_class = request.form.get("class_target")

        if note_text:
            data["notes"].append({
                "author_id": user_id,
                "author_name": user["name"],
                "author_role": user["role"],
                "text": note_text,
                "visibility": visibility,
                "target_id": target_id if visibility == "private_student" else None,
                "target_class": target_class if visibility == "class" else None,
                "target_name": (
                    next((u["name"] for u in data["users"] if u["ku_id"] == target_id), "Mentor")
                    if visibility == "private_student"
                    else (target_class if visibility == "class" else None)
                )
            })
            save_data(data)
            return redirect("/notes")

    visible_notes = []
    for n in data["notes"]:
        if (
            n.get("author_id") == user_id or
            n.get("visibility") == "public" or
            (user["role"] == "student" and n.get("visibility") == "private_student" and n.get("target_id") == user_id) or
            (user["role"] == "mentor" and n.get("visibility") == "to_mentor") or
            (user["role"] == "student" and n.get("visibility") == "class" and n.get("target_class") == user.get("class"))
        ):
            visible_notes.append(n)

    students = [u for u in data["users"] if u["role"] == "student"]

    return render_template("study_notes.html", notes=visible_notes, students=students, user=user)
# ---------- LEADERBOARD ----------
@app.route("/leaderboard")
def leaderboard():
    if "user" not in session:
        return redirect("/")
    data = load_data()
    lb_data = {}
    for r in data.get("results", []):
        q_id = str(r["quiz_id"])
        if q_id not in lb_data:
            lb_data[q_id] = {}
        s_id = r["student_id"]
        percent = (r["score"] / r["total"] * 100) if r["total"] > 0 else 0
        if s_id not in lb_data[q_id] or percent > lb_data[q_id][s_id]["percentage"]:
            lb_data[q_id][s_id] = {"name": r["student_name"], "score": r["score"], "total": r["total"], "percentage": round(percent, 1), "class": r.get("student_class", "N/A"), "timestamp": r["timestamp"]}
    processed_lb = []
    for i, quiz in enumerate(data.get("quizzes", [])):
        students_ranked = sorted(lb_data.get(str(i), {}).values(), key=lambda x: x["percentage"], reverse=True)
        processed_lb.append({"title": quiz["title"], "subject": quiz.get("subject", "General"), "rankings": students_ranked})
    return render_template("leaderboard.html", leaderboard=processed_lb, user=session["user"])

# ---------- AI STUDY NOTES ----------
@app.route("/study_notes")
def study_notes():
    if "user" not in session:
        return redirect("/")
    return render_template("ai_notes.html", user=session["user"])

@app.route("/generate-study-notes", methods=["POST"])
def generate_study_notes():
    payload = request.get_json()
    topic = payload.get("topic", "")
    style = payload.get("style", "detailed")
    if not topic:
        return jsonify({"error": "Topic is required"})
    style_map = {
        "detailed": "Write detailed, well-structured study notes with clear explanations, examples, and key concepts highlighted.",
        "bullet": "Write concise bullet-point study notes. Use clear ## headers and short bullets for each key fact.",
        "summary": "Write a short executive summary of only the most critical points a student must remember."
    }
    prompt = f"""Generate {style} study notes on: "{topic}".
{style_map.get(style, style_map['detailed'])}

Structure your notes with:
- ## Section headers
- **Bold** for key terms
- ⚠️ for important warnings or common mistakes
- Examples where helpful
- A ## Key Takeaways section at the end

Keep it student-friendly."""
    notes_content, error = call_groq_api_text(GROQ_API_KEY, prompt)
    if error:
        return jsonify({"error": error})
    return jsonify({"notes": notes_content})

# ---------- AI HELPERS ----------
def call_groq_api_text(api_key, prompt):
    try:
        response = http_requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": "llama-3.3-70b-versatile", "messages": [{"role": "user", "content": prompt}], "temperature": 0.7, "max_tokens": 4000},
            timeout=30
        )
        if response.status_code != 200:
            return None, f"API error (HTTP {response.status_code})"
        return response.json()["choices"][0]["message"]["content"].strip(), None
    except Exception as e:
        return None, str(e)

def call_groq_api(api_key, prompt):
    try:
        response = http_requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={
                "model": "llama-3.3-70b-versatile",
                "messages": [
                    {"role": "system", "content": "You are a quiz generator. Output ONLY valid JSON arrays. No markdown, no backticks."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.7, "max_tokens": 4000
            },
            timeout=30
        )
        if response.status_code == 401:
            return None, "Invalid API key."
        if response.status_code != 200:
            return None, f"Groq API error (HTTP {response.status_code})"
        content = response.json()["choices"][0]["message"]["content"].strip()
        # More robust JSON extraction: find the first '[' and last ']'
        start = content.find('[')
        end = content.rfind(']')
        if start != -1 and end != -1:
            content = content[start:end+1]
        
        return json.loads(content), None

    except Exception as e:
        return None, str(e)

@app.route("/generate-ai", methods=["POST"])
def generate_ai():
    payload = request.get_json()
    subject = payload.get("subject", "General")
    quiz_type = payload.get("type", "MCQ")
    count = payload.get("count", 5)
    difficulty = payload.get("difficulty", "Medium")
    prompt = f"Generate exactly {count} {quiz_type} quiz questions about \"{subject}\" at {difficulty} difficulty.\n\n{get_format_instruction(quiz_type)}\n\nReturn ONLY raw JSON array."
    questions, error = call_groq_api(GROQ_API_KEY, prompt)
    if error:
        return jsonify({"error": error})
    return jsonify({"questions": questions})

@app.route("/generate-from-ppt", methods=["POST"])
def generate_from_ppt():
    quiz_type = request.form.get("type", "MCQ")
    count = request.form.get("count", 5)
    difficulty = request.form.get("difficulty", "Medium")
    if "ppt" not in request.files:
        return jsonify({"error": "PPT file is required."})
    ppt_file = request.files["ppt"]
    try:
        if ppt_file.filename.lower().endswith(".ppt"):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp:
                ppt_file.save(tmp.name); tmp_path = tmp.name
            try:
                ppt_data = ppt2txt.process(tmp_path)
                text = "\n".join(ppt_data.values()) if isinstance(ppt_data, dict) else str(ppt_data)
            finally:
                if os.path.exists(tmp_path): os.remove(tmp_path)
        else:
            prs = Presentation(io.BytesIO(ppt_file.read()))
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        text = text.strip()[:15000]
        if not text:
            return jsonify({"error": "No text could be extracted."})
        prompt = f"Generate exactly {count} {quiz_type} quiz questions from this presentation:\n\n{text}\n\n{get_format_instruction(quiz_type)}\n\nReturn ONLY raw JSON array."
        questions, error = call_groq_api(GROQ_API_KEY, prompt)
        if error: return jsonify({"error": error})
        return jsonify({"questions": questions})
    except Exception as e:
        return jsonify({"error": f"PPT Error: {str(e)}"})

@app.route("/generate-from-pdf", methods=["POST"])
def generate_from_pdf():
    quiz_type = request.form.get("type", "MCQ")
    count = request.form.get("count", 5)
    difficulty = request.form.get("difficulty", "Medium")
    if "pdf" not in request.files:
        return jsonify({"error": "PDF file is required."})
    pdf_file = request.files["pdf"]
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_file.read()))
        text = "".join(page.extract_text() + "\n" for page in pdf_reader.pages)[:15000]
        prompt = f"Generate exactly {count} {quiz_type} questions from this PDF:\n\n{text}\n\n{get_format_instruction(quiz_type)}\n\nReturn ONLY raw JSON array."
        questions, error = call_groq_api(GROQ_API_KEY, prompt)
        if error: return jsonify({"error": error})
        return jsonify({"questions": questions})
    except Exception as e:
        return jsonify({"error": f"PDF Error: {str(e)}"})

def get_format_instruction(quiz_type):
    if quiz_type == "MCQ":
        return 'JSON array, each: {"question":"...","options":["A","B","C","D"],"answer":"exact correct option text"}'
    elif quiz_type == "Rapid Fire":
        return 'JSON array, each: {"question":"...","answer":"short answer"}'
    elif quiz_type == "Fill in the Blank":
        return 'JSON array, each: {"question":"sentence with _____","answer":"word"}'

    elif quiz_type == "Short Answer":
        return (
            'JSON array, each: {'
            '"question":"Clear short answer question",'
            '"answer_paragraphs":['
            '"Paragraph 1 — main definition or direct answer (2-3 sentences with key terms bolded in **bold**)",'
            '"Paragraph 2 — explanation of how/why it works (2-3 sentences)",'
            '"Paragraph 3 — formula, rule, or key equation if applicable (e.g. Program = Algorithm + Data Structure)"'
            ']}'
        )

    elif quiz_type == "Long Answer":
        return (
            'JSON array, each: {'
            '"question":"Detailed long answer question",'
            '"definition":"1 clear sentence defining the topic with key terms",'
            '"paragraphs":["Paragraph explaining concept in depth (3-4 sentences)","Second paragraph with more detail"],'
            '"types":['
            '{"name":"Type/Category 1","definition":"what it is","examples":["example1","example2"],"applications":["real use 1","real use 2"]},'
            '{"name":"Type/Category 2","definition":"what it is","examples":["example1","example2"],"applications":["real use 1","real use 2"]}'
            '],'
            '"conclusion":"2-sentence summary"}'
        )

    elif quiz_type in ("Coding", "Code"):
        return (
            'JSON array, each: {'
            '"question":"Write complete working code from scratch: [specific task with language]",'
            '"definition":"1-sentence definition of the data structure or concept",'
            '"structure_code":"struct or class definition code snippet if applicable, else empty string",'
            '"concept_points":["Key concept point 1","Key concept point 2","Key concept point 3"],'
            '"diagram_before":"text diagram before operation e.g. Head -> [10] -> [20] -> [30] -> NULL",'
            '"diagram_after":"text diagram after operation e.g. Head -> [10] -> [20] -> [30] -> [40] -> NULL",'
            '"full_code":"complete working code with inline comments"}'
        )

    return 'JSON array, each: {"question":"...","answer":"..."}'


def build_question_prompt(sec_type, difficulty, topic, material, q_num, total):
    """Build a focused prompt for a single exam question."""
    mat = material[:4000]

    if sec_type == "Short Answer":
        return f"""Generate question #{q_num} of {total} — Short Answer — {difficulty} difficulty.
Topic: {topic}
Material: {mat}

Return a SINGLE JSON object (not array) with this EXACT structure — fill ALL fields with real content:
{{
  "question": "A clear short-answer question about {topic}",
  "answer_paragraphs": [
    "Paragraph 1: Main definition with key terms highlighted (2-3 sentences, use **word** to bold important terms)",
    "Paragraph 2: How it works or why it is important (2-3 sentences with **bold** key terms)",
    "Paragraph 3: A key formula, rule, or example related to the topic"
  ]
}}

Return ONLY the JSON object. No N/A. No empty strings. Real content only."""

    elif sec_type == "Long Answer":
        return f"""Generate question #{q_num} of {total} — Long Answer — {difficulty} difficulty.
Topic: {topic}
Material: {mat}

Return a SINGLE JSON object with this EXACT structure — fill ALL fields with real content:
{{
  "question": "A detailed long-answer question about {topic}",
  "definition": "One clear sentence defining the main topic with key terms",
  "paragraphs": [
    "Explanatory paragraph 1 about the concept (3-4 sentences with **bold** key terms)",
    "Explanatory paragraph 2 with more depth (3-4 sentences with **bold** key terms)"
  ],
  "types": [
    {{
      "name": "First type or category name",
      "definition": "What this type is",
      "examples": ["example 1", "example 2", "example 3"],
      "applications": ["Real-world use 1", "Real-world use 2"]
    }},
    {{
      "name": "Second type or category name",
      "definition": "What this type is",
      "examples": ["example 1", "example 2"],
      "applications": ["Real-world use 1", "Real-world use 2"]
    }}
  ],
  "conclusion": "Two-sentence conclusion summarising the answer"
}}

Return ONLY the JSON object. No N/A. No empty strings. Real content only."""

    elif sec_type in ("Coding", "Code"):
        lang = "C++" if difficulty in ("Medium","Hard") else "C"
        return f"""Generate question #{q_num} of {total} — Coding — {difficulty} difficulty in {lang}.
Topic: {topic}
Material: {mat}

Return a SINGLE JSON object with this EXACT structure — fill ALL fields with REAL working code:
{{
  "question": "Write complete working {lang} code from scratch: [specific coding task related to {topic}]",
  "definition": "One sentence defining the data structure or algorithm used",
  "structure_code": "struct or class definition code here (use \n for newlines)",
  "concept_points": ["Step 1 of the algorithm", "Step 2", "Step 3"],
  "diagram_before": "Head -> [10] -> [20] -> [30] -> NULL",
  "diagram_after": "Head -> [10] -> [20] -> [30] -> [40] -> NULL",
  "full_code": "Complete working {lang} code with comments (use \n for newlines)"
}}

Return ONLY the JSON object. full_code must be COMPLETE working code, not a placeholder."""

    elif sec_type == "MCQ":
        return f"""Generate question #{q_num} of {total} — MCQ — {difficulty} difficulty.
Topic: {topic}
Material: {mat}

Return a SINGLE JSON object:
{{
  "question": "A clear MCQ question about {topic}",
  "options": ["Option A text", "Option B text", "Option C text", "Option D text"],
  "answer": "The exact text of the correct option"
}}

Return ONLY the JSON object."""

    else:
        return f"""Generate question #{q_num} of {total} about {topic}.
Material: {mat}
Return a single JSON object: {{"question":"...","answer":"..."}}"""


def build_batch_prompt(sec_type, difficulty, topic, material, count):
    """Build a single prompt to generate all questions for a section at once.

    Style: B.Tech student writing in an exam copy — short sentences, simple
    vocabulary, no AI-flavored phrases. Field names match what build_doc reads.
    """
    mat = material[:3000]

    student_style_rules = """
HANDWRITTEN STUDENT-STYLE RULES (apply to EVERY answer field):
- Write like a B.Tech student answering in an exam copy, NOT like a textbook or AI.
- Use SHORT, simple sentences (max ~15 words each).
- Use simple everyday English. Avoid words like: comprehensive, robust, paradigm,
  facilitate, leverage, encompasses, intricate, multifaceted, holistic, seamless,
  utilize (use "use"), demonstrate (use "show"), ascertain (use "find").
- NEVER start with: "In essence", "Furthermore", "Moreover", "Additionally",
  "It is important to note", "It should be noted", "Notably", "In conclusion".
- Each bullet/point must be ONE short line — no nested clauses.
- Use **bold** ONLY on key technical terms (1-2 per field max).
- NEVER use N/A, placeholders, or empty strings. Always write real content.
- Topic stays anchored: every answer must be about the actual topic given,
  using ONLY facts from the study material below."""

    if sec_type == "Short Answer":
        example = ('[{"question":"Define Data Structure.",'
                   '"key_points":['
                   '"A **Data Structure** is a way of storing data in computer memory.",'
                   '"It helps us access and use data faster.",'
                   '"Common types are array, stack, queue, and linked list.",'
                   '"Choosing the right one improves program speed."'
                   ']}]')
        return f"""Generate EXACTLY {count} Short Answer questions about: {topic}
Difficulty: {difficulty}
Study Material: {mat}

Return a JSON array of EXACTLY {count} objects. Each object MUST follow this EXACT format:
{example}
{student_style_rules}

ADDITIONAL RULES:
- Field name MUST be "key_points" (a list of 3-5 short strings).
- Each key_point under 20 words.
- Return ONLY the JSON array starting with [ and ending with ]."""

    elif sec_type == "Long Answer":
        example = ('[{"question":"Explain linear and non-linear data structures.",'
                   '"definition":"Data structures are ways of storing data in memory so we can use it easily.",'
                   '"key_points":['
                   '"**Linear** structures store data one after another.",'
                   '"Examples: array, stack, queue, linked list.",'
                   '"**Non-linear** structures store data in levels or links.",'
                   '"Examples: tree, graph.",'
                   '"Linear is used for marks list or undo/redo.",'
                   '"Non-linear is used for file system or social network."'
                   '],'
                   '"conclusion":"We pick the type based on the problem we need to solve."}]')
        return f"""Generate EXACTLY {count} Long Answer questions about: {topic}
Difficulty: {difficulty}
Study Material: {mat}

Return a JSON array of EXACTLY {count} objects. Each object MUST follow this EXACT format:
{example}
{student_style_rules}

ADDITIONAL RULES (Long Answer):
- "definition": 1-2 short sentences. Plain student language.
- "key_points": list of 5-7 short strings. Each under 20 words.
  Mix definition lines, examples, and applications inside the list.
- "conclusion": 1 short final sentence. Do NOT start with "In conclusion".
- Return ONLY the JSON array."""

    elif sec_type in ("Coding", "Code"):
        lang = "C++" if difficulty in ("Medium","Hard") else "C"
        example = ('[{"question":"Write complete C++ code to insert a node at end of singly linked list.",'
                   '"requirements":['
                   '"Use a struct Node with int data and Node* next.",'
                   '"Make a function to insert at end.",'
                   '"Show output by inserting 10, 20, 30."'
                   '],'
                   '"answer":"#include<iostream>\\nusing namespace std;\\nstruct Node { int data; Node* next; };\\nNode* head = NULL;\\nvoid insert(int val) {\\n    Node* newNode = new Node();\\n    newNode->data = val; newNode->next = NULL;\\n    if(head==NULL) { head=newNode; return; }\\n    Node* temp = head;\\n    while(temp->next!=NULL) temp=temp->next;\\n    temp->next = newNode;\\n}\\nint main() { insert(10); insert(20); insert(30); return 0; }"}]')
        return f"""Generate EXACTLY {count} Coding questions about: {topic}
Difficulty: {difficulty}, Language: {lang}
Study Material: {mat}

Return a JSON array of EXACTLY {count} objects. Each object MUST follow this EXACT format:
{example}
{student_style_rules}

ADDITIONAL RULES (Coding):
- Field "requirements" = list of 3-4 short instruction strings (under 15 words each).
- Field "answer" = COMPLETE working {lang} code as a single string (use \\n for newlines).
- The code MUST include #include and main() and actually compile.
- NO placeholders. Real code only.
- Return ONLY the JSON array."""

    elif sec_type == "MCQ":
        return f"""Generate EXACTLY {count} MCQ questions about: {topic}
Difficulty: {difficulty}
Study Material: {mat}

Return a JSON array of EXACTLY {count} objects:
[{{"question":"...","options":["A option","B option","C option","D option"],"answer":"exact correct option text"}}]

RULES:
- Exactly 4 options per question.
- "answer" MUST match one option text exactly.
- Question wording should be clear and short, not AI-flavored.
- Return ONLY the JSON array."""

    else:
        return f"""Generate EXACTLY {count} exam questions about {topic} from this material: {mat}
Return ONLY a JSON array: [{{"question":"...","answer":"..."}}]
Style: write the answer like a B.Tech student in an exam — short, simple sentences."""



@app.route("/generate-exam", methods=["POST"])
def generate_exam():
    if not DOCX_AVAILABLE:
        return jsonify({"error": "Exam generation requires 'python-docx' library. Please install it on the server."}), 500
    
    if "file" not in request.files:

        return jsonify({"error": "No file uploaded"}), 400
    
    ref_file = request.files["file"]
    header = json.loads(request.form.get("header", "{}"))
    sections = json.loads(request.form.get("sections", "[]"))

    # 1. Extract Text
    try:
        content = ref_file.read()
        if ref_file.filename.lower().endswith(".pdf"):
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            extracted_pages = []
            for page in pdf_reader.pages:
                t = page.extract_text()
                if t: extracted_pages.append(t + "\n")
            text = "".join(extracted_pages)[:20000]

        elif ref_file.filename.lower().endswith((".ppt", ".pptx")):
            prs = Presentation(io.BytesIO(content))
            extracted_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        extracted_text.append(shape.text)
            text = "\n".join(extracted_text)[:20000]
        else:
            return jsonify({"error": "Unsupported file format"}), 400
    except Exception as e:
        return jsonify({"error": f"Extraction Error: {str(e)}"}), 500


    if not text.strip():
        return jsonify({"error": "No readable text found in file"}), 400

    # 2. Generate Sections via AI
    exam_data = []
    for sec in sections:
        sec_name  = sec.get('name', 'Section')
        sec_type  = sec.get('type', 'MCQ')
        sec_count = int(sec.get('count', 5))
        sec_diff  = sec.get('difficulty', 'Medium')
        sec_marks = sec.get('marks', 2)

        questions = []
        last_error = ""
        for attempt in range(3):
            try:
                if attempt > 0:
                    time.sleep(attempt * 8)
                # Use the actual subject as topic, not the section name ("Section A" etc.)
                actual_topic = (header.get("subject") or "").strip() or sec_name
                batch_prompt = build_batch_prompt(sec_type, sec_diff, actual_topic, text, sec_count)
                raw_resp = http_requests.post(
                    "https://api.groq.com/openai/v1/chat/completions",
                    headers={"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"},
                    json={
                        "model": "llama-3.1-8b-instant",
                        "messages": [
                            {"role": "system", "content": "You are a university exam question generator. Return ONLY a valid JSON array. No markdown. No backticks. Never use N/A."},
                            {"role": "user", "content": batch_prompt}
                        ],
                        "temperature": 0.3,
                        "max_tokens": 2048
                    },
                    timeout=60
                )
                if raw_resp.status_code == 200:
                    raw = raw_resp.json()["choices"][0]["message"]["content"].strip()
                    raw = raw.replace("```json","").replace("```","").strip()
                    start = raw.find("["); end = raw.rfind("]") + 1
                    if start != -1 and end > 0:
                        questions = json.loads(raw[start:end])
                        if isinstance(questions, list) and len(questions) > 0:
                            break
                elif raw_resp.status_code == 429:
                    last_error = "Rate limited"
                    continue
            except Exception as e:
                last_error = str(e)
                continue

        if not questions:
            try:
                time.sleep(5)
                fmt = get_format_instruction(sec_type)
                fb_prompt = f"Generate {sec_count} {sec_type} exam questions about {sec_name}. Format: {fmt}. Return ONLY JSON array."
                questions, _ = call_groq_api(GROQ_API_KEY, fb_prompt)
                if not isinstance(questions, list):
                    questions = []
            except Exception as e2:
                last_error += f" | {str(e2)}"
                questions = []

        if questions:
            exam_data.append({"name": sec_name, "type": sec_type, "marks": sec_marks, "questions": questions})
            print(f"[EXAM GEN] Section '{sec_name}' OK: {len(questions)} questions")
        else:
            print(f"[EXAM GEN] Section '{sec_name}' failed: {last_error}")
        time.sleep(5)


    if not exam_data:
        return jsonify({"error": "AI failed to generate questions for any section."}), 500

    # 3. Create Documents
    def create_doc(is_answer_key=False):
        doc = Document()
        
        # Header
        inst = header.get('institution', 'INSTITUTION NAME').upper()
        h1 = doc.add_heading(inst, 0)
        h1.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        details = doc.add_paragraph()
        details.alignment = WD_ALIGN_PARAGRAPH.CENTER
        details.add_run(f"{header.get('title', 'Examination')}\n").bold = True
        details.add_run(f"Subject: {header.get('subject', 'General')}\n")
        details.add_run(f"Course: {header.get('degree', 'N/A')} | Date: {header.get('date', 'TBD')}\n")
        details.add_run(f"Time: {header.get('time', '3 Hours')} | Total Marks: {header.get('totalMarks', '100')}")

        if not is_answer_key:
            seat = doc.add_paragraph()
            seat.add_run("\nStudent Name / Seat No: ___________________________").italic = True
        else:
            key_title = doc.add_paragraph()
            key_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            key_title.add_run("\n*** TEACHER'S REFERENCE / ANSWER KEY ***").bold = True

        doc.add_paragraph("-" * 40)

        # Content
        for sec in exam_data:
            s_title = doc.add_heading(sec['name'], level=1)
            s_info = doc.add_paragraph()
            s_info.add_run(f"[{sec['type']} - {len(sec['questions'])} Questions x {sec['marks']} Marks each]").italic = True
            
            for i, q in enumerate(sec['questions']):
                q_p = doc.add_paragraph()
                q_p.add_run(f"Q{i+1}. {q.get('question', 'Question text missing')} ").bold = True
                q_p.add_run(f"({sec['marks']} Marks)")
                
                # MCQ options
                if sec['type'] == 'MCQ' and 'options' in q:
                    for j, opt in enumerate(q['options']):
                        doc.add_paragraph(f"   {chr(65+j)}) {opt}")

                if is_answer_key:
                    # ── SHORT ANSWER: paragraphs like a model answer ──────────
                    if sec['type'] == 'Short Answer':
                        ans_p = doc.add_paragraph()
                        ans_p.add_run("Answer:").bold = True
                        # Support both new (answer_paragraphs) and old (key_points/answer) format
                        paras = q.get('answer_paragraphs', [])
                        if not paras:
                            # Fallback: use key_points or answer
                            kps = q.get('key_points', [])
                            if kps:
                                paras = kps
                            elif q.get('answer') and q['answer'] != 'N/A':
                                paras = [q['answer']]
                        for para in paras:
                            p = doc.add_paragraph()
                            p.paragraph_format.left_indent = Inches(0.3)
                            parts = str(para).split('**')
                            for idx2, part in enumerate(parts):
                                run = p.add_run(part)
                                run.bold = (idx2 % 2 == 1)

                    # ── LONG ANSWER: Definition + Types + Applications ────────
                    elif sec['type'] == 'Long Answer':
                        defn = q.get('definition', '')
                        if defn:
                            dp = doc.add_paragraph()
                            dp.add_run(defn)

                        for para in q.get('paragraphs', []):
                            p = doc.add_paragraph()
                            p.paragraph_format.left_indent = Inches(0.1)
                            parts = str(para).split('**')
                            for idx2, part in enumerate(parts):
                                run = p.add_run(part)
                                run.bold = (idx2 % 2 == 1)

                        types = q.get('types', [])
                        # Fallback: use key_points as bullet list if no types
                        if not types and q.get('key_points'):
                            kp_h = doc.add_paragraph()
                            kp_h.add_run("Key Points:").bold = True
                            for pt in q['key_points']:
                                bp = doc.add_paragraph(f"  • {pt}")
                                bp.paragraph_format.left_indent = Inches(0.3)
                        elif types:
                            th = doc.add_paragraph()
                            th.add_run("Types:").bold = True
                            for t in types:
                                tn = doc.add_paragraph()
                                tn.paragraph_format.left_indent = Inches(0.3)
                                tn.add_run(f"{t.get('name','')}: ").bold = True
                                tn.add_run(t.get('definition',''))
                                exs = t.get('examples', [])
                                if exs:
                                    ep = doc.add_paragraph(f"      Examples: {', '.join(exs)}")
                                    ep.paragraph_format.left_indent = Inches(0.5)
                                apps = t.get('applications', [])
                                for app in apps:
                                    ap2 = doc.add_paragraph(f"      • {app}")
                                    ap2.paragraph_format.left_indent = Inches(0.5)

                        conc = q.get('conclusion', '')
                        if conc:
                            cp = doc.add_paragraph()
                            cp.add_run("Conclusion: ").bold = True
                            cp.add_run(conc)

                    # ── CODING: Def + Structure + Diagram + Full Code ─────────
                    elif sec['type'] in ('Coding', 'Code'):
                        defn = q.get('definition', '')
                        if defn:
                            dp = doc.add_paragraph()
                            dp.add_run("Definition: ").bold = True
                            dp.add_run(defn)

                        struct = q.get('structure_code', '')
                        if struct:
                            doc.add_paragraph().add_run("Node Structure:").bold = True
                            sb = doc.add_paragraph(struct)
                            if sb.runs:
                                sb.runs[0].font.name = 'Courier New'
                                sb.runs[0].font.size = Pt(9)
                            sb.paragraph_format.left_indent = Inches(0.4)

                        pts = q.get('concept_points', [])
                        if pts:
                            ph = doc.add_paragraph()
                            ph.add_run("Key Concepts:").bold = True
                            for pt in pts:
                                bp = doc.add_paragraph(f"  • {pt}")
                                bp.paragraph_format.left_indent = Inches(0.3)

                        db = q.get('diagram_before', '')
                        da = q.get('diagram_after', '')
                        if db or da:
                            doc.add_paragraph().add_run("Diagram:").bold = True
                            if db:
                                b1 = doc.add_paragraph()
                                b1.add_run("Before: ").bold = True
                                b1.add_run(db)
                            if da:
                                b2 = doc.add_paragraph()
                                b2.add_run("After:  ").bold = True
                                b2.add_run(da)

                        code = q.get('full_code', q.get('answer', ''))
                        if code:
                            doc.add_paragraph().add_run("Solution Code:").bold = True
                            cb = doc.add_paragraph(code)
                            if cb.runs:
                                cb.runs[0].font.name = 'Courier New'
                                cb.runs[0].font.size = Pt(9)
                            cb.paragraph_format.left_indent = Inches(0.4)

                    else:
                        ans_p = doc.add_paragraph()
                        ans_p.add_run(f"Answer: {q.get('answer', 'N/A')}").italic = True

                else:
                    # ── QUESTION PAPER — NO answer spaces at all ──────────────
                    pass  # MCQ options already printed; no lines for other types

                doc.add_paragraph()  # Spacer

        return doc

    # 4. Save to static file and return URL
    mode = request.form.get("mode", "paper")
    sub_clean = "".join(c for c in header.get('subject', 'Subject') if c.isalnum() or c in (' ', '-', '_')).strip()
    
    gen_dir = os.path.join(app.static_folder, "generated")
    if not os.path.exists(gen_dir):
        os.makedirs(gen_dir, exist_ok=True)
        
    # Periodic cleanup (older than 10 mins)
    now = datetime.datetime.now()
    for f in os.listdir(gen_dir):
        fpath = os.path.join(gen_dir, f)
        if os.path.getmtime(fpath) < (now.timestamp() - 600):
            try: os.remove(fpath)
            except: pass

    if mode == "key":
        doc = create_doc(is_answer_key=True)
        filename = f"AnswerKey_{sub_clean}_{uuid.uuid4().hex[:8]}.docx"
    else:
        doc = create_doc(is_answer_key=False)
        filename = f"QuestionPaper_{sub_clean}_{uuid.uuid4().hex[:8]}.docx"

    file_path = os.path.join(gen_dir, filename)
    doc.save(file_path)
    
    return jsonify({
        "success": True, 
        "url": url_for('static', filename=f"generated/{filename}"),
        "filename": filename
    })



# ---------- ANSWER CHECKER ----------

UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), 'static', 'answer_pdfs')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)


@app.route("/check-answers-page")
def check_answers_page():
    if "user" not in session:
        return redirect("/")
    data = load_data()
    user = session["user"]
    my_reports = sorted(
        [r for r in data.get("answer_reports", []) if r.get("student_id") == user["ku_id"]],
        key=lambda x: x.get("submitted_at", ""), reverse=True
    )
    return render_template("answer_checker.html", user=user, my_reports=my_reports)



# ── NLP helpers (loaded once at startup) ──────────────────────────────────────
_nlp_model = None

def get_nlp_model():
    """Lazy-load sentence-transformers model. Falls back gracefully if not installed."""
    global _nlp_model
    if _nlp_model is not None:
        return _nlp_model
    try:
        from sentence_transformers import SentenceTransformer, util as st_util
        _nlp_model = SentenceTransformer("all-MiniLM-L6-v2")
        return _nlp_model
    except Exception:
        return None


def nlp_semantic_similarity(text_a: str, text_b: str) -> float:
    """Return cosine similarity 0.0–1.0 between two texts using sentence-transformers."""
    model = get_nlp_model()
    if model is None or not text_a.strip() or not text_b.strip():
        return -1.0   # -1 signals NLP unavailable
    try:
        from sentence_transformers import util as st_util
        import torch
        emb_a = model.encode(text_a, convert_to_tensor=True)
        emb_b = model.encode(text_b, convert_to_tensor=True)
        score = float(st_util.cos_sim(emb_a, emb_b)[0][0])
        return round(score, 4)
    except Exception:
        return -1.0


def nlp_keyword_coverage(student_text: str, reference_text: str) -> dict:
    """
    Extract key noun/verb tokens from reference, check how many student covered.
    Returns {"covered": [...], "missing": [...], "coverage_pct": float}
    Uses spacy if available, else simple word overlap.
    """
    try:
        import spacy
        try:
            nlp = spacy.load("en_core_web_sm")
        except OSError:
            raise ImportError("spacy model not downloaded")

        ref_doc = nlp(reference_text.lower())
        stu_doc = nlp(student_text.lower())

        # Key terms: nouns, proper nouns, verbs (lemmatised, length > 3)
        ref_keywords = set(
            tok.lemma_ for tok in ref_doc
            if tok.pos_ in ("NOUN", "PROPN", "VERB") and len(tok.lemma_) > 3
            and not tok.is_stop and tok.is_alpha
        )
        stu_lemmas = set(tok.lemma_ for tok in stu_doc if tok.is_alpha)

        covered = sorted(ref_keywords & stu_lemmas)
        missing = sorted(ref_keywords - stu_lemmas)
        pct = round(len(covered) / max(len(ref_keywords), 1) * 100, 1)
        return {"covered": covered[:10], "missing": missing[:10], "coverage_pct": pct}

    except Exception:
        # Fallback: simple word overlap
        ref_words = set(w.lower() for w in reference_text.split() if len(w) > 3)
        stu_words = set(w.lower() for w in student_text.split() if len(w) > 3)
        covered = sorted(ref_words & stu_words)
        missing = sorted(ref_words - stu_words)
        pct = round(len(covered) / max(len(ref_words), 1) * 100, 1)
        return {"covered": covered[:10], "missing": missing[:10], "coverage_pct": pct}


def nlp_analyse_student_answer(student_answer: str, model_answer: str, max_marks: float) -> dict:
    """
    Run full NLP analysis on one answer pair.
    Returns dict with similarity, keyword_coverage, suggested_marks_pct, nlp_available.
    """
    similarity   = nlp_semantic_similarity(student_answer, model_answer)
    kw_coverage  = nlp_keyword_coverage(student_answer, model_answer)
    nlp_available = similarity >= 0

    if nlp_available:
        # Weighted score: 60% semantic similarity + 40% keyword coverage
        kw_score   = kw_coverage["coverage_pct"] / 100.0
        nlp_score  = 0.6 * similarity + 0.4 * kw_score
        suggested_pct = round(max(0.0, min(1.0, nlp_score)) * 100, 1)
    else:
        suggested_pct = None

    return {
        "nlp_available":   nlp_available,
        "semantic_similarity": round(similarity, 4) if nlp_available else None,
        "keyword_coverage":    kw_coverage,
        "nlp_suggested_pct":   suggested_pct,
    }


@app.route("/check-answers", methods=["POST"])
def check_answers():
    if "user" not in session:
        return jsonify({"error": "Not logged in"}), 401

    # ── 1. Validate: PDF only ─────────────────────────────────────────────────
    if "pdf" not in request.files or request.files["pdf"].filename == "":
        return jsonify({"error": "No file uploaded. Please select a PDF."}), 400

    pdf_file = request.files["pdf"]
    filename = pdf_file.filename.lower()

    ALLOWED_EXTENSIONS = (".pdf", ".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp")
    if not any(filename.endswith(ext) for ext in ALLOWED_EXTENSIONS):
        return jsonify({"error": "Unsupported file type. Upload a PDF or a photo (JPG/PNG) of your answer sheet."}), 400
    
    is_image = any(filename.endswith(ext) for ext in (".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"))

    user        = session["user"]
    subject     = request.form.get("subject",    "General").strip()
    total_marks = int(request.form.get("totalMarks", 30))
    context     = request.form.get("context",    "").strip()
    strictness  = request.form.get("strictness", "standard")
    exam_name   = request.form.get("examName",   "").strip() or subject

    # ── 2. Save PDF to disk ───────────────────────────────────────────────────
    report_id = str(uuid.uuid4())
    safe_name = f"{report_id}.pdf"
    pdf_path  = os.path.join(UPLOAD_FOLDER, safe_name)
    pdf_bytes = pdf_file.read()

    try:
        with open(pdf_path, "wb") as f:
            f.write(pdf_bytes)
    except Exception as e:
        return jsonify({"error": f"Could not save PDF: {str(e)}"}), 500

    pdf_url = f"/static/answer_pdfs/{safe_name}"

    # ── 3. Extract text (PDF or raw image) — Gemini Vision first, EasyOCR fallback ──
    # Strategy:
    #   1. For digital PDFs, try fast PyPDF2 path first
    #   2. If little/no text, use Gemini Vision API (works on handwriting + full pages)
    #   3. If Gemini fails (quota, network), fall back to EasyOCR (offline)
    ocr_method = "unknown"
    extracted_text = ""

    # Save the file to disk so OCR modules can read by path
    # (We already saved pdf_bytes to pdf_path above as part of the upload)

    # Step 1: For digital PDFs only, try PyPDF2 first
    if not is_image:
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_bytes))
            quick_text = "".join((p.extract_text() or "") + "\n\n"
                                 for p in pdf_reader.pages).strip()
            if len(quick_text) > 100:
                extracted_text = quick_text
                ocr_method = "pypdf2"
                print(f"  PyPDF2 extracted {len(quick_text)} chars (digital PDF)")
        except Exception as e:
            print(f"  PyPDF2 failed (will try OCR): {e}")

    # Step 2: If no usable text yet, try Gemini Vision
    if not extracted_text:
        try:
            from gemini_extract import extract_text as gemini_extract_text
            extracted_text = gemini_extract_text(pdf_path)
            ocr_method = "gemini"
            print(f"  Gemini extracted {len(extracted_text)} chars")
            print(f"  Preview: {extracted_text[:200]}")
        except Exception as gemini_error:
            print(f"  Gemini failed, falling back to EasyOCR: {gemini_error}")
            try:
                from ocr_extract import extract_text as ocr_extract_text
                extracted_text = ocr_extract_text(pdf_path)
                ocr_method = "easyocr_fallback"
                print(f"  EasyOCR extracted {len(extracted_text)} chars")
            except Exception as ocr_error:
                print(f"  Both Gemini and EasyOCR failed: {ocr_error}")
                return jsonify({
                    "error": f"Could not read your file. Try a clearer photo/scan, "
                             f"make sure handwriting is dark on a white background, "
                             f"and avoid shadows or glare. (OCR error: {str(ocr_error)})"
                }), 400

    if not extracted_text.strip():
        return jsonify({
            "error": "No readable text found. Try a clearer photo/scan, "
                     "make sure handwriting is dark on a white background, "
                     "and avoid shadows or glare."
        }), 400

    # ── 4. Pre-NLP: run sentence-transformer similarity if context provided ───
    #    We compare the full student text against the exam context/questions
    #    to get a high-level semantic score before per-question LLM grading.
    pre_nlp = {}
    if context:
        pre_nlp = nlp_analyse_student_answer(extracted_text, context, total_marks)

    # ── 5. Build NLP-enriched grading prompt ─────────────────────────────────
    strictness_map = {
        "lenient":  "Be generous — reward partial understanding. Minor spelling/phrasing errors should NOT cost marks.",
        "standard": "Grade fairly like an experienced teacher. Give partial credit for partially correct answers.",
        "strict":   "Grade strictly — full marks only for complete, precise, well-structured answers.",
    }

    # Build NLP context block for the prompt
    nlp_context_block = ""
    if pre_nlp.get("nlp_available"):
        sim  = pre_nlp.get("semantic_similarity", 0)
        kw   = pre_nlp.get("keyword_coverage", {})
        sugg = pre_nlp.get("nlp_suggested_pct", 0)
        nlp_context_block = f"""
NLP PRE-ANALYSIS (use as grading guidance, not hard constraint):
- Semantic similarity between student paper and exam questions: {sim:.2%}
- Key concepts COVERED by student: {", ".join(kw.get("covered", [])) or "none detected"}
- Key concepts MISSING from student answer: {", ".join(kw.get("missing", [])) or "none"}
- Keyword coverage: {kw.get("coverage_pct", 0)}%
- NLP suggested overall score range: ~{sugg}% (use as a cross-check, your assessment takes priority)
"""

    ctx_block = f"\n\nORIGINAL EXAM QUESTIONS / CONTEXT:\n{context}" if context else ""

    prompt = f"""You are an expert NLP-powered university examiner grading a student answer sheet for: "{subject}".
Total marks: {total_marks}
Grading style: {strictness_map.get(strictness, strictness_map["standard"])}
{ctx_block}
{nlp_context_block}

STUDENT ANSWER SHEET (extracted from PDF):
---
{extracted_text[:12000]}
---

Instructions:
1. Identify EVERY question answered (Q1, Q2, Q1a, Q1b, 1., 2. etc.)
2. For each question, perform deep NLP-style evaluation:
   - Check semantic correctness (does the answer MEAN the right thing, even if phrased differently?)
   - Check keyword/concept coverage (are key terms present?)
   - Check completeness (is the answer fully developed?)
   - Check clarity (is the answer structured and clear?)
3. Assign marks using 0.5 increments
4. Write a 2-sentence holistic verdict like a professor
5. List exactly what the student got right (specific points)
6. List exactly what is missing or wrong (specific gaps)
7. Give one concrete improvement tip per question

Return ONLY valid JSON (no markdown, no backticks):
{{
  "scored": <integer or .5 total earned>,
  "overall_feedback": "<2-3 sentence NLP-informed professor-style overall assessment>",
  "strongest_area": "<section/topic student did best at>",
  "weakest_area": "<section/topic needing most work>",
  "encouragement": "<1 motivating sentence>",
  "nlp_insights": "<1-2 sentences about semantic quality of the overall paper>",
  "questions": [
    {{
      "question_text": "<e.g. Q1a: Cost Function>",
      "max_marks": <integer>,
      "earned_marks": <number, 0.5 increments>,
      "student_answer": "<verbatim or close summary of what student wrote>",
      "holistic_verdict": "<2-sentence NLP-informed professor comment>",
      "semantic_quality": "<Excellent|Good|Partial|Poor> — how semantically correct is the answer>",
      "correct_points": ["<specific concept/point student got right>"],
      "gaps": ["<specific missing concept or error — be very precise>"],
      "suggestion": "<one concrete improvement tip>"
    }}
  ]
}}
Rules:
- scored = sum of all earned_marks, must NOT exceed {total_marks}
- correct_points empty ONLY if zero marks; gaps empty ONLY if full marks
- semantic_quality must be one of: Excellent, Good, Partial, Poor
- Be specific in gaps: say "Missing time complexity O(V+E) for BFS" not "Incomplete"
- Paraphrased correct answers SHOULD get credit (semantic understanding > exact wording)"""

    # ── 6. Call Groq LLM ─────────────────────────────────────────────────────
    try:
        response = http_requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"},
            json={
                "model": "llama-3.3-70b-versatile",
                "messages": [
                    {"role": "system", "content": "You are a strict JSON-only NLP-powered university exam grader. Output ONLY valid JSON. No markdown. No text outside the JSON."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.2,
                "max_tokens": 2048
            },
            timeout=60
        )
        if response.status_code != 200:
            return jsonify({"error": f"AI service error (HTTP {response.status_code})"}), 500

        raw   = response.json()["choices"][0]["message"]["content"].strip()
        start = raw.find("{"); end = raw.rfind("}") + 1
        if start == -1 or end == 0:
            return jsonify({"error": "AI returned unexpected response. Please try again."}), 500

        result = json.loads(raw[start:end])
        if result.get("scored", 0) > total_marks:
            result["scored"] = total_marks

    except json.JSONDecodeError:
        return jsonify({"error": "AI returned malformed response. Please try again."}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    # ── 7. Run per-question NLP analysis and attach to each question ──────────
    questions = result.get("questions", [])
    model_answers_from_context = []
    if context:
        # Try to extract model answers from context (if mentor pasted questions with marks)
        import re
        # Simple extraction: lines that look like answers
        lines = context.split("\n")
        model_answers_from_context = [l.strip() for l in lines if len(l.strip()) > 20]

    nlp_available_global = get_nlp_model() is not None

    for i, q in enumerate(questions):
        student_ans = q.get("student_answer", "")
        # Use model answer hint from context if available, else use correct_points as reference
        model_ref = ""
        if i < len(model_answers_from_context):
            model_ref = model_answers_from_context[i]
        elif q.get("correct_points"):
            model_ref = " ".join(q["correct_points"])

        if model_ref and student_ans:
            q_nlp = nlp_analyse_student_answer(student_ans, model_ref, q.get("max_marks", 1))
            q["nlp_analysis"] = {
                "semantic_similarity": q_nlp.get("semantic_similarity"),
                "keyword_coverage_pct": q_nlp["keyword_coverage"].get("coverage_pct"),
                "missing_keywords": q_nlp["keyword_coverage"].get("missing", [])[:5],
                "nlp_suggested_pct": q_nlp.get("nlp_suggested_pct"),
                "nlp_available": q_nlp.get("nlp_available", False)
            }
        else:
            q["nlp_analysis"] = {"nlp_available": False}

    result["questions"] = questions
    result["pre_nlp"]   = pre_nlp
    result["nlp_available"] = nlp_available_global

    # ── 8. Save report ────────────────────────────────────────────────────────
    data = load_data()
    report = {
        "id":            report_id,
        "student_id":    user["ku_id"],
        "student_name":  user["name"],
        "student_class": user.get("class", "N/A"),
        "exam_name":     exam_name,
        "subject":       subject,
        "total_marks":   total_marks,
        "scored":        result.get("scored", 0),
        "strictness":    strictness,
        "submitted_at":  datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "pdf_url":       pdf_url,
        "status":        "pending",
        "mentor_note":   "",
        "nlp_available": nlp_available_global,
        "ocr_method":    ocr_method,
        "result":        result
    }
    data["answer_reports"].append(report)
    save_data(data)

    result["report_id"]  = report_id
    result["pdf_url"]    = pdf_url
    result["ocr_method"] = ocr_method
    return jsonify(result)


@app.route("/delete-answer-report/<report_id>", methods=["POST"])
def delete_answer_report(report_id):
    if "user" not in session or session["user"].get("role") != "mentor":
        return jsonify({"error": "Unauthorised"}), 403
    data = load_data()
    # Also delete the stored PDF
    report = next((r for r in data.get("answer_reports", []) if r["id"] == report_id), None)
    if report and report.get("pdf_url"):
        try:
            pdf_path = os.path.join(os.path.dirname(__file__), "static", "answer_pdfs", f"{report_id}.pdf")
            if os.path.exists(pdf_path):
                os.remove(pdf_path)
        except Exception:
            pass
    data["answer_reports"] = [r for r in data.get("answer_reports", []) if r.get("id") != report_id]
    save_data(data)
    return redirect("/mentor")


@app.route("/generate-exam-both", methods=["POST"])
def generate_exam_both():
    """Generate both question paper and answer key as a ZIP containing two .docx files."""
    if not DOCX_AVAILABLE:
        return jsonify({"error": "Exam generation requires python-docx library."}), 500

    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    ref_file = request.files["file"]
    header   = json.loads(request.form.get("header", "{}"))
    sections = json.loads(request.form.get("sections", "[]"))

    # 1. Extract text
    try:
        content = ref_file.read()
        if ref_file.filename.lower().endswith(".pdf"):
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            text = "".join((p.extract_text() or "") + "\n" for p in pdf_reader.pages)[:20000]
        elif ref_file.filename.lower().endswith((".ppt", ".pptx")):
            prs = Presentation(io.BytesIO(content))
            text = "\n".join(
                shape.text for slide in prs.slides
                for shape in slide.shapes if hasattr(shape, "text")
            )[:20000]
        else:
            return jsonify({"error": "Unsupported file format. Use PDF or PPTX."}), 400
    except Exception as e:
        return jsonify({"error": f"Extraction Error: {str(e)}"}), 500

    if not text.strip():
        return jsonify({"error": "No readable text found in the uploaded file."}), 400

    # 2. Generate questions via AI
    exam_data = []
    for sec in sections:
        sec_name  = sec.get("name", "Section")
        sec_type  = sec.get("type", "MCQ")
        sec_count = int(sec.get("count", 5))
        sec_diff  = sec.get("difficulty", "Medium")
        sec_marks = sec.get("marks", 2)

        questions = []
        last_error2 = ""
        
        # Try with retry + exponential backoff
        for attempt in range(3):
            try:
                if attempt > 0:
                    wait = attempt * 8
                    print(f"[EXAM BOTH] Retry {attempt} for '{sec_name}' after {wait}s...")
                    time.sleep(wait)
                
                # Use the actual subject as topic, not the section name ("Section A" etc.)
                actual_topic = (header.get("subject") or "").strip() or sec_name
                batch_prompt = build_batch_prompt(sec_type, sec_diff, actual_topic, text, sec_count)
                raw_resp = http_requests.post(
                    "https://api.groq.com/openai/v1/chat/completions",
                    headers={"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"},
                    json={
                        "model": "llama-3.1-8b-instant",
                        "messages": [
                            {"role": "system", "content": "You are a university exam question generator. Return ONLY a valid JSON array. No markdown. No backticks. Never use N/A."},
                            {"role": "user", "content": batch_prompt}
                        ],
                        "temperature": 0.3,
                        "max_tokens": 2048
                    },
                    timeout=60
                )
                if raw_resp.status_code == 200:
                    raw = raw_resp.json()["choices"][0]["message"]["content"].strip()
                    raw = raw.replace("```json","").replace("```","").strip()
                    start = raw.find("["); end = raw.rfind("]") + 1
                    if start != -1 and end > 0:
                        questions = json.loads(raw[start:end])
                        if isinstance(questions, list) and len(questions) > 0:
                            break  # Success!
                elif raw_resp.status_code == 429:
                    last_error2 = "Rate limited"
                    continue
                else:
                    last_error2 = f"API {raw_resp.status_code}"
                    break
            except Exception as e:
                last_error2 = str(e)
                continue

        # Final fallback with simplest possible prompt
        if not questions:
            try:
                time.sleep(5)
                fmt = get_format_instruction(sec_type)
                fb = f"Generate {sec_count} {sec_type} exam questions about {sec_name}. Format: {fmt}. Return ONLY JSON array."
                questions, _ = call_groq_api(GROQ_API_KEY, fb)
                if not isinstance(questions, list):
                    questions = []
            except Exception as e2:
                last_error2 += f" | {str(e2)}"
                questions = []

        if questions:
            exam_data.append({"name": sec_name, "type": sec_type,
                               "marks": sec_marks, "questions": questions})
            print(f"[EXAM BOTH] Section '{sec_name}' OK: {len(questions)} questions")
        else:
            print(f"[EXAM BOTH] Section '{sec_name}' failed: {last_error2}")
        
        time.sleep(5)  # Wait between sections

    if not exam_data:
        return jsonify({"error": "AI failed to generate questions. Check Flask terminal for details. Try with a smaller PDF or fewer sections."}), 500

    # 3. Build document helper
    def build_doc(is_answer_key):
        doc = Document()
        for section in doc.sections:
            section.top_margin    = Inches(1)
            section.bottom_margin = Inches(1)
            section.left_margin   = Inches(1)
            section.right_margin  = Inches(1)

        h1 = doc.add_heading(header.get("institution", "INSTITUTION").upper(), 0)
        h1.alignment = WD_ALIGN_PARAGRAPH.CENTER

        meta = doc.add_paragraph()
        meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
        meta.add_run(f"{header.get('title', 'Examination')}\n").bold = True
        meta.add_run(f"Subject: {header.get('subject', '')}   |   Course: {header.get('degree', '')}\n")
        meta.add_run(f"Date: {header.get('date', 'TBD')}   |   Time: {header.get('time', '3 Hours')}   |   Total Marks: {header.get('totalMarks', '100')}")

        if is_answer_key:
            kp = doc.add_paragraph()
            kp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            kr = kp.add_run("★  ANSWER KEY — TEACHER'S COPY (CONFIDENTIAL)  ★")
            kr.bold = True
        else:
            sp = doc.add_paragraph()
            sp.add_run("Student Name: _________________________   Seat No: ___________").italic = True

        doc.add_paragraph()
        q_serial = 1
        for sec in exam_data:
            doc.add_heading(sec["name"], level=1)
            info = doc.add_paragraph()
            info.add_run(f"[{sec['type']}  ·  {len(sec['questions'])} Questions  ·  {sec['marks']} Mark(s) each]").italic = True
            for q in sec["questions"]:
                qp = doc.add_paragraph()
                qp.add_run(f"Q{q_serial}.  {q.get('question', '')}").bold = True
                qp.add_run(f"  [{sec['marks']} M]")
                if sec["type"] == "MCQ" and "options" in q:
                    for j, opt in enumerate(q["options"]):
                        doc.add_paragraph(f"   ({chr(65+j)})  {opt}")

                if is_answer_key:
                    if sec["type"] == "Short Answer":
                        kpts = q.get("key_points") or []
                        if kpts:
                            ap = doc.add_paragraph(); ap.add_run("   Model Answer (Key Points):").bold = True
                            for pt in kpts:
                                if pt and str(pt).strip().upper() != "N/A":
                                    doc.add_paragraph(f"   • {pt}")
                        else:
                            # Fallback: try single-string answer if no key_points
                            fb = q.get("answer", "")
                            if fb and str(fb).strip().upper() != "N/A":
                                ap = doc.add_paragraph(); ap.add_run("   Answer: ").bold = True
                                ap.add_run(str(fb))
                    elif sec["type"] == "Long Answer":
                        defn = q.get("definition", "").strip()
                        if defn and defn.upper() != "N/A":
                            dp = doc.add_paragraph(); dp.add_run("   Definition: ").bold = True; dp.add_run(defn)
                        kpts = q.get("key_points") or []
                        clean_kpts = [pt for pt in kpts if pt and str(pt).strip().upper() != "N/A"]
                        if clean_kpts:
                            kpp = doc.add_paragraph(); kpp.add_run("   Key Points:").bold = True
                            for pt in clean_kpts:
                                doc.add_paragraph(f"   • {pt}")
                        conc = q.get("conclusion", "").strip()
                        if conc and conc.upper() != "N/A":
                            cp = doc.add_paragraph(); cp.add_run("   Conclusion: ").bold = True; cp.add_run(conc)
                    elif sec["type"] in ("Coding","Code"):
                        reqs = [r for r in (q.get("requirements") or []) if r and str(r).strip().upper() != "N/A"]
                        if reqs:
                            rp = doc.add_paragraph(); rp.add_run("   Requirements:").bold = True
                            for r in reqs: doc.add_paragraph(f"   • {r}")
                        code = (q.get("answer") or q.get("full_code") or "").strip()
                        if code and code.upper() != "N/A":
                            cp = doc.add_paragraph(); cp.add_run("   Solution Code:").bold = True
                            cb = doc.add_paragraph(code)
                            if cb.runs: cb.runs[0].font.name = "Courier New"; cb.runs[0].font.size = Pt(9)
                    else:
                        ans = q.get("answer", "").strip()
                        if ans and ans.upper() != "N/A":
                            ap = doc.add_paragraph()
                            ap.add_run(f"   Answer: {ans}").italic = True
                # Question paper — NO answer lines, just questions
                doc.add_paragraph()
                q_serial += 1
        return doc

    # 4. Save both to BytesIO and ZIP
    sub_clean = "".join(c for c in header.get("subject", "Exam")
                        if c.isalnum() or c in (" ", "-", "_")).strip().replace(" ", "_")
    uid = uuid.uuid4().hex[:8]

    paper_buf = io.BytesIO()
    key_buf   = io.BytesIO()
    build_doc(False).save(paper_buf)
    build_doc(True ).save(key_buf)
    paper_buf.seek(0); key_buf.seek(0)

    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"QuestionPaper_{sub_clean}_{uid}.docx", paper_buf.read())
        zf.writestr(f"AnswerKey_{sub_clean}_{uid}.docx",     key_buf.read())
    zip_buf.seek(0)

    return send_file(
        zip_buf,
        mimetype="application/zip",
        as_attachment=True,
        download_name=f"ExamPack_{sub_clean}_{uid}.zip"
    )


@app.errorhandler(404)
def page_not_found(e):
    return render_template("404.html"), 404

if __name__ == "__main__":
    app.run(debug=True)